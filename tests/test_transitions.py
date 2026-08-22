from __future__ import annotations

import json
import shutil
import zipfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.oxml.ns import qn

from ppt_agent import service
from ppt_agent.errors import PptAgentError
from ppt_agent.ooxml import MC_NS, restore_transition_options, strip_transition_elements, transitions_by_slide
from ppt_agent.paths import document_id, revision

FIXTURE = Path("fixtures/synthetic/synthetic.pptx")


def _base_deck(path: Path) -> Path:
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(path)
    return path


def _add_direct_transition(path: Path, slide_index: int = 0, attrs: dict | None = None) -> None:
    prs = Presentation(path)
    sld = prs.slides[slide_index]._element
    tr = sld.makeelement(qn("p:transition"), attrs or {})
    tr.append(sld.makeelement(qn("p:fade"), {}))
    sld.append(tr)
    prs.save(path)


def _add_wrapped_transition(path: Path, slide_index: int = 0, attrs: dict | None = None) -> None:
    prs = Presentation(path)
    sld = prs.slides[slide_index]._element
    tr = sld.makeelement(qn("p:transition"), attrs or {})
    tr.append(sld.makeelement(qn("p:fade"), {}))
    fallback = sld.makeelement("{%s}Fallback" % MC_NS)
    fallback.append(tr)
    alternate = sld.makeelement("{%s}AlternateContent" % MC_NS)
    alternate.append(fallback)
    sld.append(alternate)
    prs.save(path)


def _slide_xml(path: Path) -> str:
    with zipfile.ZipFile(path) as archive:
        return archive.read("ppt/slides/slide1.xml").decode("utf-8")


def test_transitions_by_slide_reads_direct_transition(tmp_path: Path) -> None:
    deck = _base_deck(tmp_path / "direct.pptx")
    _add_direct_transition(deck, attrs={"spd": "med"})
    assert transitions_by_slide(deck)[0] == {"type": "fade", "speed": "med"}


def test_transitions_by_slide_reads_mc_wrapped_transition(tmp_path: Path) -> None:
    deck = _base_deck(tmp_path / "wrapped.pptx")
    _add_wrapped_transition(deck, attrs={"spd": "med"})
    assert transitions_by_slide(deck)[0] == {"type": "fade", "speed": "med"}


def test_transitions_by_slide_reports_none_without_transition(tmp_path: Path) -> None:
    deck = _base_deck(tmp_path / "plain.pptx")
    assert transitions_by_slide(deck)[0] is None


def test_strip_transition_elements_removes_direct_and_wrapped(tmp_path: Path) -> None:
    deck = _base_deck(tmp_path / "both.pptx")
    _add_direct_transition(deck)
    _add_wrapped_transition(deck)
    strip_transition_elements(deck, {0})
    xml = _slide_xml(deck)
    assert "transition" not in xml
    assert "AlternateContent" not in xml


def test_inspect_json_reports_wrapped_transition(tmp_path: Path) -> None:
    from ppt_agent.engine import inspect_json

    deck = tmp_path / "deck.pptx"
    shutil.copy2(FIXTURE, deck)
    _add_wrapped_transition(deck, slide_index=0, attrs={"spd": "med"})
    data = inspect_json(deck)
    assert data["slides"]["0"]["_transition"] == {"type": "fade", "speed": "med"}


def test_structural_diff_reports_wrapped_transition_change(tmp_path: Path) -> None:
    from ppt_agent.engine import structural_diff

    before = tmp_path / "before.pptx"
    after = tmp_path / "after.pptx"
    _base_deck(before)
    shutil.copy2(before, after)
    _add_wrapped_transition(before, attrs={"spd": "med"})

    changes = structural_diff(before, after)
    assert "~ transition speed=med type=fade -> none" in changes


def test_structural_diff_ignores_identical_wrapped_transitions(tmp_path: Path) -> None:
    from ppt_agent.engine import structural_diff

    before = tmp_path / "before.pptx"
    after = tmp_path / "after.pptx"
    _base_deck(before)
    shutil.copy2(before, after)
    _add_wrapped_transition(before, attrs={"spd": "med"})
    _add_wrapped_transition(after, attrs={"spd": "med"})

    changes = structural_diff(before, after)
    assert "transition" not in changes


def test_set_slide_replaces_wrapped_transition(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    source = tmp_path / "deck.pptx"
    shutil.copy2(FIXTURE, source)
    _add_wrapped_transition(source, slide_index=0, attrs={"spd": "med"})
    monkeypatch.setenv("PPT_AGENT_STATE_DIR", str(tmp_path / "state"))
    monkeypatch.setattr(service.wps, "finalize", lambda path: {"wps_version": "test"})
    monkeypatch.setattr(service.wps, "reopen_verify", lambda path: {"wps_version": "test"})
    patch = {
        "document_id": document_id(source),
        "revision": revision(source),
        "operations": [{"op": "set_slide", "slide": 0, "transition": {"type": "push", "dir": "l"}}],
    }
    patch_path = tmp_path / "patch.json"
    patch_path.write_text(json.dumps(patch), encoding="utf-8")
    response = service.apply(source, patch_path, allow_risk={"external_relationship"})
    assert response["ok"] is True

    candidate = tmp_path / "deck.agent.candidate.pptx"
    xml = _slide_xml(candidate)
    assert "AlternateContent" not in xml
    assert '<p:push dir="l"/>' in xml or "<p:push" in xml
    assert transitions_by_slide(candidate)[0] == {"type": "push", "dir": "l"}


def test_set_slide_none_removes_wrapped_transition(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    source = tmp_path / "deck.pptx"
    shutil.copy2(FIXTURE, source)
    _add_wrapped_transition(source, slide_index=0, attrs={"spd": "med"})
    monkeypatch.setenv("PPT_AGENT_STATE_DIR", str(tmp_path / "state"))
    monkeypatch.setattr(service.wps, "finalize", lambda path: {"wps_version": "test"})
    patch = {
        "document_id": document_id(source),
        "revision": revision(source),
        "operations": [{"op": "set_slide", "slide": 0, "transition": "none"}],
    }
    patch_path = tmp_path / "patch.json"
    patch_path.write_text(json.dumps(patch), encoding="utf-8")
    response = service.apply(source, patch_path, allow_risk={"external_relationship"})
    assert response["ok"] is True

    candidate = tmp_path / "deck.agent.candidate.pptx"
    xml = _slide_xml(candidate)
    assert "transition" not in xml
    assert "AlternateContent" not in xml
    assert transitions_by_slide(candidate)[0] is None


def test_wps_normalized_transition_options_are_restored_before_publish(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "deck.pptx"
    shutil.copy2(FIXTURE, source)
    monkeypatch.setenv("PPT_AGENT_STATE_DIR", str(tmp_path / "state"))

    def normalizing_finalize(path: Path) -> dict:
        # Simulate the observed WPS save behavior: keep type and timing,
        # drop spd and effect direction attributes.
        prs = Presentation(path)
        sld = prs.slides[0]._element
        tr = sld.find(qn("p:transition"))
        assert tr is not None
        tr.attrib.pop("spd", None)
        for child in tr:
            child.attrib.pop("dir", None)
            child.attrib.pop("orient", None)
            break
        prs.save(path)
        return {"wps_version": "test"}

    monkeypatch.setattr(service.wps, "finalize", normalizing_finalize)
    monkeypatch.setattr(service.wps, "reopen_verify", lambda path: {"wps_version": "test"})
    patch = {
        "document_id": document_id(source),
        "revision": revision(source),
        "operations": [{
            "op": "set_slide",
            "slide": 0,
            "transition": {"type": "push", "dir": "l", "speed": "fast"},
        }],
    }
    patch_path = tmp_path / "patch.json"
    patch_path.write_text(json.dumps(patch), encoding="utf-8")
    response = service.apply(source, patch_path, allow_risk={"external_relationship"})
    assert response["ok"] is True

    candidate = tmp_path / "deck.agent.candidate.pptx"
    assert transitions_by_slide(candidate)[0] == {"type": "push", "dir": "l", "speed": "fast"}


# uncover/pull fidelity --------------------------------------------------------

def _add_pull_transition(path: Path, slide_index: int = 0) -> None:
    prs = Presentation(path)
    sld = prs.slides[slide_index]._element
    tr = sld.makeelement(qn("p:transition"), {})
    tr.append(sld.makeelement(qn("p:pull"), {}))
    sld.append(tr)
    prs.save(path)


def test_restore_uncover_pull_dir_and_speed(tmp_path: Path) -> None:
    deck = _base_deck(tmp_path / "pull.pptx")
    _add_pull_transition(deck)
    restore_transition_options(deck, {0: {"type": "uncover", "dir": "l", "speed": "fast"}})
    assert transitions_by_slide(deck)[0] == {"type": "uncover", "dir": "l", "speed": "fast"}


def test_restore_wrapped_uncover_pull_dir_and_speed(tmp_path: Path) -> None:
    deck = _base_deck(tmp_path / "pull-wrapped.pptx")
    prs = Presentation(deck)
    sld = prs.slides[0]._element
    tr = sld.makeelement(qn("p:transition"), {})
    tr.append(sld.makeelement(qn("p:pull"), {}))
    fallback = sld.makeelement("{%s}Fallback" % MC_NS)
    fallback.append(tr)
    alternate = sld.makeelement("{%s}AlternateContent" % MC_NS)
    alternate.append(fallback)
    sld.append(alternate)
    prs.save(deck)
    restore_transition_options(deck, {0: {"type": "uncover", "dir": "ru", "speed": "slow"}})
    assert transitions_by_slide(deck)[0] == {"type": "uncover", "dir": "ru", "speed": "slow"}


def test_restore_mismatched_transition_type_raises_fidelity_failure(tmp_path: Path) -> None:
    deck = _base_deck(tmp_path / "push.pptx")
    _add_direct_transition(deck, attrs={})  # writes p:fade
    with pytest.raises(PptAgentError) as caught:
        restore_transition_options(deck, {0: {"type": "push", "dir": "l"}})
    assert caught.value.code == "TRANSITION_FIDELITY_FAILED"


# surgical ZIP edit ------------------------------------------------------------

def test_restore_changes_only_target_slide_zip_entry(tmp_path: Path) -> None:
    deck = tmp_path / "deck.pptx"
    shutil.copy2(FIXTURE, deck)
    _add_pull_transition(deck, slide_index=0)
    _add_pull_transition(deck, slide_index=1)

    with zipfile.ZipFile(deck) as archive:
        before = {info.filename: archive.read(info.filename) for info in archive.infolist()}

    restore_transition_options(deck, {0: {"type": "uncover", "dir": "l", "speed": "fast"}})

    with zipfile.ZipFile(deck) as archive:
        after = {info.filename: archive.read(info.filename) for info in archive.infolist()}

    changed = [name for name in before if before[name] != after.get(name)]
    assert changed == ["ppt/slides/slide1.xml"]
    assert set(before) == set(after)
    assert transitions_by_slide(deck)[0] == {"type": "uncover", "dir": "l", "speed": "fast"}
    # slide 2 untouched
    assert transitions_by_slide(deck)[1] == {"type": "uncover"}


def test_replace_with_retry_recovers_from_transient_winerror(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    from ppt_agent.ooxml import _replace_with_retry

    calls = {"count": 0}
    real_replace = __import__("os").replace

    def flaky(source, target):
        calls["count"] += 1
        if calls["count"] == 1:
            error = OSError("transient busy")
            error.winerror = 5
            raise error
        return real_replace(source, target)

    monkeypatch.setattr("ppt_agent.ooxml.os.replace", flaky)
    a = tmp_path / "a"
    b = tmp_path / "b"
    a.write_text("x", encoding="utf-8")
    b.write_text("y", encoding="utf-8")
    _replace_with_retry(a, b)
    assert b.read_text(encoding="utf-8") == "x"
    assert calls["count"] == 2


def test_replace_with_retry_gives_up_after_bounded_attempts(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    from ppt_agent.ooxml import _replace_with_retry

    calls = {"count": 0}

    def always_busy(source, target):
        calls["count"] += 1
        error = OSError("always busy")
        error.winerror = 32
        raise error

    monkeypatch.setattr("ppt_agent.ooxml.os.replace", always_busy)
    a = tmp_path / "a"
    b = tmp_path / "b"
    a.write_text("x", encoding="utf-8")
    b.write_text("y", encoding="utf-8")
    with pytest.raises(OSError):
        _replace_with_retry(a, b)
    assert calls["count"] == 5
    assert b.read_text(encoding="utf-8") == "y"  # 目标未改变


# diff correctness -------------------------------------------------------------

def test_diff_equivalent_direct_and_wrapped_transition_reports_no_change(tmp_path: Path) -> None:
    from ppt_agent.engine import structural_diff

    before = tmp_path / "before.pptx"
    after = tmp_path / "after.pptx"
    _base_deck(before)
    shutil.copy2(before, after)
    _add_direct_transition(before, attrs={"spd": "med"})
    _add_wrapped_transition(after, attrs={"spd": "med"})

    changes = structural_diff(before, after)
    assert "transition" not in changes
    assert changes == "No structural differences."


def test_diff_keeps_plain_text_containing_transition_word(tmp_path: Path) -> None:
    from ppt_agent.engine import structural_diff

    before = tmp_path / "before.pptx"
    after = tmp_path / "after.pptx"
    _base_deck(before)
    prs = Presentation(before)
    prs.slides[0].shapes.add_textbox(0, 0, 4000000, 500000).text = "transition old"
    prs.save(before)
    shutil.copy2(before, after)
    prs = Presentation(after)
    prs.slides[0].shapes[0].text_frame.text = "transition new"
    prs.save(after)
    _add_direct_transition(after, attrs={})

    changes = structural_diff(before, after)
    assert "transition old -> transition new" in changes
    assert changes.count("~ transition ") == 1
    assert "~ transition none -> type=fade" in changes


def test_diff_reports_single_line_for_wrapped_to_direct_change(tmp_path: Path) -> None:
    from ppt_agent.engine import structural_diff

    before = tmp_path / "before.pptx"
    after = tmp_path / "after.pptx"
    _base_deck(before)
    _add_wrapped_transition(before, attrs={"spd": "med"})
    shutil.copy2(before, after)
    prs = Presentation(after)
    sld = prs.slides[0]._element
    for alternate in list(sld.findall("{%s}AlternateContent" % MC_NS)):
        sld.remove(alternate)
    tr = sld.makeelement(qn("p:transition"), {})
    push = sld.makeelement(qn("p:push"), {"dir": "l"})
    tr.append(push)
    sld.append(tr)
    prs.save(after)

    changes = structural_diff(before, after)
    assert changes.count("~ transition ") == 1
    assert "~ transition speed=med type=fade -> dir=l type=push" in changes
