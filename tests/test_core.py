from __future__ import annotations

import hashlib
import json
import subprocess
import zipfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from ppt_agent.errors import PptAgentError
from ppt_agent.models import PatchRequest
from ppt_agent.ooxml import security_scan, validate_pptx
from ppt_agent.paths import accepted_path, candidate_path, document_id, revision, source_path
from ppt_agent import service


@pytest.fixture
def deck(tmp_path: Path) -> Path:
    path = tmp_path / "deck.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text = "Original title"
    prs.save(path)
    return path


@pytest.fixture(autouse=True)
def state_dir(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.setenv("PPT_AGENT_STATE_DIR", str(tmp_path / "state"))


def test_paths_keep_original_separate(deck: Path) -> None:
    assert candidate_path(deck).name == "deck.agent.candidate.pptx"
    assert accepted_path(deck).name == "deck.agent.pptx"
    assert source_path(candidate_path(deck)) == deck
    assert source_path(accepted_path(deck)) == deck
    assert document_id(deck) == document_id(deck)


def test_validate_and_scan_plain_deck(deck: Path) -> None:
    validate_pptx(deck)
    report = security_scan(deck)
    assert report["safe"] is True
    assert report["counts"]["slides"] == 1


def test_resize_scale_and_line_add_shape_translate_without_internal_error() -> None:
    from ppt_agent.engine import translate_operation

    assert translate_operation({"op": "resize", "object": "s0:s2", "scale": 0.8}) == {
        "op": "resize", "slide": 0, "shape": "s2", "scale": 0.8,
    }
    assert translate_operation({"op": "resize", "object": "s0:s2", "width": 4.0, "height": 1.5}) == {
        "op": "resize", "slide": 0, "shape": "s2", "size": [4.0, 1.5],
    }
    assert translate_operation({
        "op": "add_shape", "slide": 0, "kind": "line", "from": [0.0, 0.0], "to": [5.0, 5.0],
    }) == {
        "op": "add-shape", "slide": 0, "kind": "line", "from": [0.0, 0.0], "to": [5.0, 5.0],
    }
    assert translate_operation({
        "op": "add_shape", "slide": 0, "kind": "textbox", "x": 1.0, "y": 2.0, "width": 4.0, "height": 1.5,
    }) == {
        "op": "add-shape", "slide": 0, "kind": "textbox", "at": [1.0, 2.0], "size": [4.0, 1.5],
    }


def test_native_chart_workbook_is_not_ole_risk() -> None:
    fixture = Path("fixtures/synthetic/synthetic.pptx")
    report = security_scan(fixture)
    assert "ole" not in {item["type"] for item in report["risks"]}
    assert report["counts"]["charts"] == 3


def test_patch_schema_requires_expect_count_for_selector(deck: Path) -> None:
    with pytest.raises(ValueError):
        PatchRequest.model_validate({
            "request_id": "x",
            "document_id": document_id(deck),
            "revision": revision(deck),
            "operations": [{"op": "set_text", "selector": {"role": "title"}, "text": "x"}],
        })


def test_patch_schema_is_discriminated_and_rejects_unknown_fields(deck: Path) -> None:
    with pytest.raises(ValueError):
        PatchRequest.model_validate({
            "request_id": "x",
            "document_id": document_id(deck),
            "revision": revision(deck),
            "operations": [{"op": "move", "object": "s0:s2", "x": 1, "y": 2, "mystery": True}],
        })


def test_focused_inspect_returns_only_requested_fields(deck: Path) -> None:
    data = service.inspect(deck, purpose="layout", no_state=True)["data"]["structure"]
    shape = data["slides"]["0"]["s2"]

    assert data["fields"] == ["geometry", "identity", "name", "text", "type"]
    assert set(shape) == {"identity", "name", "type", "pos", "size", "paragraphs"}
    assert len(json.dumps(data, ensure_ascii=False).encode("utf-8")) < 64 * 1024


def test_apply_is_atomic_and_idempotent(deck: Path, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.setattr(service.wps, "finalize", lambda path: {"wps_version": "test"})
    original_hash = hashlib.sha256(deck.read_bytes()).hexdigest()
    request = {
        "request_id": "req-1",
        "document_id": document_id(deck),
        "revision": revision(deck),
        "operations": [{"op": "set_text", "object": "s0:s2", "text": "Changed"}],
    }
    patch = tmp_path / "patch.json"
    patch.write_text(json.dumps(request), encoding="utf-8")

    first = service.apply(deck, patch)
    second = service.apply(deck, patch)

    assert first == second
    assert candidate_path(deck).exists()
    assert hashlib.sha256(deck.read_bytes()).hexdigest() == original_hash
    assert "Changed" in service.inspect(candidate_path(deck), no_state=True)["data"]["structure"]


def test_apply_rejects_revision_conflict_before_wps(
    deck: Path, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    wps_started = False

    def fail_if_started(*args, **kwargs):
        nonlocal wps_started
        wps_started = True
        raise AssertionError("WPS must not start")

    monkeypatch.setattr(service.wps, "finalize", fail_if_started)
    monkeypatch.setattr(service.wps, "apply_wps_operations", fail_if_started)
    patch = tmp_path / "patch.json"
    patch.write_text(json.dumps({
        "request_id": "req-conflict",
        "document_id": document_id(deck),
        "revision": "sha256:" + "0" * 64,
        "operations": [{"op": "set_text", "object": "s0:s2", "text": "Changed"}],
    }), encoding="utf-8")
    with pytest.raises(PptAgentError) as caught:
        service.apply(deck, patch)
    assert caught.value.code == "REVISION_CONFLICT"
    assert caught.value.details["document_unchanged"] is True
    assert wps_started is False
    assert not candidate_path(deck).exists()


def test_revision_conflict_points_to_existing_candidate_and_corrected_argv(
    deck: Path, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    candidate = candidate_path(deck)
    presentation = Presentation(deck)
    presentation.slides[0].shapes[0].text = "Candidate title"
    presentation.save(candidate)
    patch = tmp_path / "candidate-bound.json"
    patch.write_text(json.dumps({
        "request_id": "req-candidate-conflict",
        "document_id": document_id(deck),
        "revision": revision(deck),
        "operations": [{"op": "set_text", "object": "s0:s2", "text": "Changed"}],
    }), encoding="utf-8")
    monkeypatch.setattr(
        service.wps,
        "finalize",
        lambda *args, **kwargs: (_ for _ in ()).throw(AssertionError("WPS must not start")),
    )

    with pytest.raises(PptAgentError) as caught:
        service.apply(deck, patch)

    error = caught.value
    assert error.code == "REVISION_CONFLICT"
    assert error.details["corrected_path"] == str(candidate.resolve())
    assert error.details["corrected_argv"] == [
        "apply",
        str(candidate.resolve()),
        "--patch",
        str(patch),
    ]
    assert error.details["reinspect_argv"] == ["inspect", str(candidate.resolve()), "--for", "edit"]
    assert error.details["revision_update_required"] == revision(candidate)
    assert error.details["patch_revision_matches_corrected_path"] is False
    assert error.details["document_unchanged"] is True


def test_candidate_bound_patch_conflict_returns_immediately_usable_candidate_argv(
    deck: Path, tmp_path: Path
) -> None:
    candidate = candidate_path(deck)
    presentation = Presentation(deck)
    presentation.slides[0].shapes[0].text = "Candidate title"
    presentation.save(candidate)
    patch = tmp_path / "candidate-bound.json"
    patch.write_text(json.dumps({
        "request_id": "req-candidate-bound",
        "document_id": document_id(deck),
        "revision": revision(candidate),
        "operations": [{"op": "set_text", "object": "s0:s2", "text": "Changed"}],
    }), encoding="utf-8")

    with pytest.raises(PptAgentError) as caught:
        service.apply(deck, patch, restart=True)

    details = caught.value.details
    assert caught.value.code == "REVISION_CONFLICT"
    assert details["corrected_path"] == str(candidate.resolve())
    assert details["corrected_argv"] == [
        "apply",
        str(candidate.resolve()),
        "--patch",
        str(patch),
    ]
    assert details["patch_revision_matches_corrected_path"] is True
    assert details["revision_update_required"] is None
    assert details["document_unchanged"] is True


def test_bad_operation_does_not_publish_candidate(deck: Path, tmp_path: Path) -> None:
    patch = tmp_path / "patch.json"
    patch.write_text(json.dumps({
        "request_id": "req-bad",
        "document_id": document_id(deck),
        "revision": revision(deck),
        "operations": [{"op": "motion_path_animation", "object": "s0:s2"}],
    }), encoding="utf-8")
    with pytest.raises(PptAgentError) as caught:
        service.apply(deck, patch)
    assert caught.value.code == "UNSUPPORTED_OPERATION"
    assert not candidate_path(deck).exists()


def test_qa_detects_out_of_bounds(deck: Path) -> None:
    prs = Presentation(deck)
    prs.slides[0].shapes.add_textbox(Inches(12), Inches(7), Inches(3), Inches(2)).text = "outside"
    prs.save(deck)
    report = service.qa(deck)["data"]
    assert report["error_count"] >= 1
    assert any(item["code"] == "OUT_OF_BOUNDS" for item in report["issues"])


def test_qa_suggests_an_executable_patch_without_modifying(deck: Path) -> None:
    original = deck.read_bytes()
    prs = Presentation(deck)
    shape = prs.slides[0].shapes[0]
    shape.left = Inches(12)
    shape.top = Inches(7)
    shape.text_frame.paragraphs[0].runs[0].font.size = Pt(8)
    prs.save(deck)
    before_qa = deck.read_bytes()

    report = service.qa(deck, "presentation", suggest_fixes=True)["data"]
    patch = PatchRequest.model_validate(report["suggested_patch"])

    assert report["suggested_fix_count"] == 2
    assert {operation.op for operation in patch.operations} == {"move", "set_style"}
    assert deck.read_bytes() == before_qa
    assert before_qa != original


def test_invalid_zip_is_rejected(tmp_path: Path) -> None:
    path = tmp_path / "bad.pptx"
    path.write_bytes(b"not a zip")
    with pytest.raises(PptAgentError) as caught:
        validate_pptx(path)
    assert caught.value.code == "INVALID_PPTX"


def test_external_relationship_is_reported() -> None:
    report = security_scan(Path("fixtures/synthetic/synthetic.pptx"))
    assert "external_relationship" in {item["type"] for item in report["risks"]}


def test_diff_reports_before_to_after_direction(deck: Path, tmp_path: Path) -> None:
    after = tmp_path / "after.pptx"
    prs = Presentation(deck)
    prs.slides[0].shapes[0].text = "After title"
    prs.save(after)
    changes = service.diff(deck, after)["data"]["changes"]
    assert "Original title -> After title" in changes


def test_render_removes_stale_page_images(
    deck: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    doc_id = document_id(deck)
    store = service.StateStore(doc_id)
    render_id = revision(deck).split(":", 1)[1][:16]
    target = store.render_dir / render_id
    target.mkdir(parents=True)
    (target / "slide-2.jpg").write_bytes(b"stale")

    def fake_export(path: Path, pdf_path: Path) -> dict:
        pdf_path.write_bytes(b"pdf")
        return {"wps_version": "test"}

    def fake_run(command, **kwargs):
        Path(command[-1] + "-1.jpg").write_bytes(b"fresh")
        return subprocess.CompletedProcess(command, 0, "", "")

    monkeypatch.setattr(service.wps, "export_pdf", fake_export)
    monkeypatch.setattr(service, "_pdftoppm", lambda: Path("pdftoppm.exe"))
    monkeypatch.setattr(service.subprocess, "run", fake_run)
    monkeypatch.setattr(service, "run_qa", lambda path, profile: {
        "profile": profile,
        "issue_count": 0,
        "error_count": 0,
        "warning_count": 0,
        "issues": [],
    })
    response = service.render(deck, pages="0")
    images = response["data"]["images"]
    assert len(images) == 1
    assert images[0].endswith("slide-1.jpg")
