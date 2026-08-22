from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from ppt_agent import service
from ppt_agent import cli
from ppt_agent.creation import create_presentation
from ppt_agent.models import CreateSpec


def test_quote_layout_preserves_every_body_item(tmp_path: Path) -> None:
    spec = CreateSpec.model_validate({
        "slides": [{"layout": "quote", "title": "Gate", "body": ["one", "two", "three"]}]
    })
    output = tmp_path / "quote.pptx"
    create_presentation(spec, output)
    text = "\n".join(shape.text for shape in Presentation(output).slides[0].shapes if shape.has_text_frame)
    assert all(item in text for item in spec.slides[0].body)


def test_existing_deck_workflow_advertises_image_replacement() -> None:
    assert "swap_image" in service.STANDARD_OPS


def test_paragraphs_all_calls_add_effect_once_at_wps_level_one(monkeypatch: pytest.MonkeyPatch, tmp_path: Path) -> None:
    captured: list[tuple] = []

    class Sequence:
        def AddEffect(self, *args):
            captured.append(args)
            return type("Effect", (), {"Timing": type("Timing", (), {})()})()

    shape = type("Shape", (), {"Id": 2})()
    slide = type("Slide", (), {
        "Shapes": [shape],
        "TimeLine": type("Timeline", (), {"MainSequence": Sequence()})(),
    })()
    presentation = type("Presentation", (), {
        "Slides": lambda self, index: slide,
        "Save": lambda self: None,
        "Close": lambda self: None,
    })()
    app = type("App", (), {"Presentations": type("Presentations", (), {"Open": lambda self, *args: presentation})()})()

    from contextlib import contextmanager
    from ppt_agent import wps

    @contextmanager
    def fake_application():
        yield app

    monkeypatch.setattr(wps, "wps_application", fake_application)
    monkeypatch.setattr(wps, "reopen_verify", lambda path: {"wps_version": "test", "slide_count": 1})
    wps.apply_wps_operations(tmp_path / "deck.pptx", [{
        "op": "add_animation",
        "object": "s0:s2",
        "effect": "appear",
        "trigger": "with_previous",
        "paragraphs": "all",
    }])
    assert len(captured) == 1
    assert captured[0][2] == 1


def test_wps_reopen_starts_after_writer_application_exits(
    monkeypatch: pytest.MonkeyPatch, tmp_path: Path
) -> None:
    from contextlib import contextmanager
    from ppt_agent import wps

    depth = 0
    reopen_depths: list[int] = []
    transition = type("Transition", (), {"EntryEffect": 0})()
    slide = type("Slide", (), {"SlideShowTransition": transition})()
    presentation = type(
        "Presentation",
        (),
        {
            "Slides": lambda self, index: slide,
            "Save": lambda self: None,
            "Close": lambda self: None,
        },
    )()
    app = type(
        "App",
        (),
        {"Presentations": type("Presentations", (), {"Open": lambda self, *args: presentation})()},
    )()

    @contextmanager
    def fake_application():
        nonlocal depth
        depth += 1
        try:
            yield app
        finally:
            depth -= 1

    def fake_reopen(path: Path) -> dict:
        reopen_depths.append(depth)
        return {"wps_version": "test", "slide_count": 1}

    monkeypatch.setattr(wps, "wps_application", fake_application)
    monkeypatch.setattr(wps, "reopen_verify", fake_reopen)
    wps.apply_wps_operations(
        tmp_path / "deck.pptx",
        [{"op": "set_transition", "slide": 0, "transition": "fade"}],
    )

    assert reopen_depths == [0]


def test_cli_usage_error_is_structured_json(capsys: pytest.CaptureFixture[str]) -> None:
    with pytest.raises(SystemExit) as caught:
        cli.main(["inspect"])
    assert caught.value.code == 1
    assert __import__("json").loads(capsys.readouterr().out)["ok"] is False
