from __future__ import annotations

import hashlib
import json
import shutil
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from ppt_agent import cli, service, wps
from ppt_agent.errors import PptAgentError
from ppt_agent.paths import accepted_path, candidate_path, document_id, revision
from ppt_agent.state import StateStore


@pytest.fixture
def deck(tmp_path: Path) -> Path:
    path = tmp_path / "deck.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text = "Original"
    prs.save(path)
    return path


@pytest.fixture(autouse=True)
def state_dir(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.setenv("PPT_AGENT_STATE_DIR", str(tmp_path / "state"))


def _patch(path: Path, request_id: str, operations: list[dict]) -> dict:
    return {
        "request_id": request_id,
        "document_id": document_id(path),
        "revision": revision(path),
        "operations": operations,
    }


def test_same_request_id_with_different_payload_is_rejected(
    deck: Path, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    monkeypatch.setattr(service.wps, "finalize", lambda path: {"wps_version": "test"})
    patch_path = tmp_path / "patch.json"
    first = _patch(deck, "same-id", [{"op": "set_text", "object": "s0:s2", "text": "One"}])
    patch_path.write_text(json.dumps(first), encoding="utf-8")
    service.apply(deck, patch_path)

    second = dict(first)
    second["operations"] = [{"op": "set_text", "object": "s0:s2", "text": "Two"}]
    patch_path.write_text(json.dumps(second), encoding="utf-8")
    with pytest.raises(PptAgentError) as caught:
        service.apply(deck, patch_path)
    assert caught.value.code == "REQUEST_ID_CONFLICT"


def test_apply_accepts_utf8_bom_and_generates_stable_request_id(
    deck: Path, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    def fake_apply(source: Path, operations: list[dict], output: Path) -> dict:
        shutil.copy2(source, output)
        return {"engine_output": "prepared"}

    monkeypatch.setattr(service, "apply_operations", fake_apply)
    monkeypatch.setattr(service.wps, "finalize", lambda path: {"wps_version": "test"})
    patch = {
        "document_id": document_id(deck),
        "revision": revision(deck),
        "operations": [{"op": "set_text", "object": "s0:s2", "text": "Changed"}],
    }
    patch_path = tmp_path / "bom-patch.json"
    patch_path.write_text(json.dumps(patch), encoding="utf-8-sig")

    first = service.apply(deck, patch_path)
    second = service.apply(deck, patch_path)
    request_ids = list(StateStore(document_id(deck)).requests())

    assert first == second
    assert len(request_ids) == 1
    assert request_ids[0].startswith("auto-")


def test_apply_accepts_complete_patch_request_from_stdin_text(
    deck: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    calls = {"engine": 0, "wps": 0}

    def fake_apply(source: Path, operations: list[dict], output: Path) -> dict:
        calls["engine"] += 1
        shutil.copy2(source, output)
        return {"engine_output": "prepared"}

    def fake_finalize(path: Path) -> dict:
        calls["wps"] += 1
        return {"wps_version": "test"}

    monkeypatch.setattr(service, "apply_operations", fake_apply)
    monkeypatch.setattr(service.wps, "finalize", fake_finalize)
    patch = {
        "document_id": document_id(deck),
        "revision": revision(deck),
        "operations": [{"op": "set_text", "object": "s0:s2", "text": "Changed"}],
    }

    result = service.apply(deck, None, patch_text="\ufeff" + json.dumps(patch))

    assert result["ok"] is True
    assert calls == {"engine": 1, "wps": 1}


def test_stdin_operations_array_is_rejected_before_wps(deck: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    wps_started = False

    def fail_if_started(*args, **kwargs):
        nonlocal wps_started
        wps_started = True
        raise AssertionError("WPS must not start")

    monkeypatch.setattr(service.wps, "finalize", fail_if_started)
    monkeypatch.setattr(service.wps, "apply_wps_operations", fail_if_started)

    with pytest.raises(PptAgentError) as caught:
        service.apply(deck, None, patch_text=json.dumps([{"op": "delete", "object": "s0:s2"}]))

    assert caught.value.code == "INVALID_PATCH"
    assert caught.value.details["document_unchanged"] is True
    assert wps_started is False
    assert not candidate_path(deck).exists()


def test_operations_array_returns_copyable_patch_request_example(deck: Path, tmp_path: Path) -> None:
    operations = [{"op": "set_text", "object": "s0:s2", "text": "Changed"}]
    patch_path = tmp_path / "operations-only.json"
    patch_path.write_text(json.dumps(operations), encoding="utf-8")

    with pytest.raises(PptAgentError) as caught:
        service.apply(deck, patch_path)

    error = caught.value
    example = error.details["corrected_example"]
    assert error.code == "INVALID_PATCH"
    assert error.next_action == "wrap_patch_request"
    assert example["document_id"] == document_id(deck)
    assert example["revision"] == revision(deck)
    assert example["operations"] == operations
    assert example["request_id"].startswith("auto-")
    assert error.details["document_unchanged"] is True


def test_invalid_json_returns_corrected_example_without_changing_document(deck: Path, tmp_path: Path) -> None:
    patch_path = tmp_path / "invalid.json"
    patch_path.write_text("{invalid", encoding="utf-8")

    with pytest.raises(PptAgentError) as caught:
        service.apply(deck, patch_path)

    details = caught.value.details
    assert caught.value.code == "INVALID_PATCH"
    assert details["corrected_example"]["document_id"] == document_id(deck)
    assert details["corrected_example"]["revision"] == revision(deck)
    assert details["document_unchanged"] is True
    assert not candidate_path(deck).exists()


@pytest.mark.parametrize(
    ("fields", "expected_groups"),
    [
        (["slide_index", "shape_id", "object_id"], set()),
        (["x", "y", "width", "height"], {"geometry"}),
        (["slide_no", "shape_id", "shape_type", "text"], {"text", "type"}),
        (["slide_index", "shape_name", "kind", "left", "top", "width", "height"], {"geometry", "name", "type"}),
        (["slides.index", "slides.shapes.id", "slides.shapes.name", "slides.shapes.x", "slides.shapes.text"], {"geometry", "name", "text"}),
        (["shape.*"], {"geometry", "name", "text", "type"}),
        (["slide_id", "paragraph_index", "animations"], {"animation"}),
        (["effect", "trigger", "paragraphs", "duration", "delay", "transition"], {"animation"}),
        (["animation.effect", "animation.trigger", "animation.paragraphs"], {"animation"}),
        (["slide.number", "slide.title", "shape.name", "shape.text", "shape.x", "shape.y", "shape.width", "shape.height"], {"geometry", "name", "text"}),
    ],
)
def test_inspect_common_field_aliases_succeed_first_time(fields: list[str], expected_groups: set[str]) -> None:
    projected = service._project_inspection(
        {
            "slide_count": 1,
            "slide_size": [13.33, 7.5],
            "slides": {
                "0": {
                    "s2": {
                        "type": "TEXT_BOX",
                        "name": "Title",
                        "pos": [1.0, 2.0],
                        "size": [3.0, 4.0],
                        "paragraphs": [{"text": "Hello"}],
                        "animations": [],
                    }
                }
            },
        },
        purpose=None,
        fields=fields,
    )

    shape = projected["slides"]["0"]["s2"]
    assert shape["identity"] == {"slide_index": 0, "shape_id": "s2", "object_id": "s0:s2"}
    assert expected_groups.issubset(set(projected["fields"]))


def test_inspect_purpose_ignores_redundant_understandable_fields() -> None:
    projected = service._project_inspection(
        {"slide_count": 1, "slide_size": [13.33, 7.5], "slides": {}},
        purpose="layout",
        fields=["slide_index", "slides.shapes.x", "shape.unknown_extra"],
    )

    assert projected["fields"] == ["geometry", "identity", "name", "text", "type"]
    assert projected["ignored_fields"] == ["shape.unknown_extra"]


@pytest.mark.parametrize("purpose", ["text-edit", "content", "edit", "layout", "animation"])
def test_focused_inspect_returns_revision_bound_apply_contract(deck: Path, purpose: str) -> None:
    payload = service.inspect(deck, purpose=purpose, no_state=True)

    template = payload["data"]["patch_template"]
    assert template == {
        "request_id": None,
        "document_id": document_id(deck),
        "revision": revision(deck),
        "operations": [],
    }
    assert payload["data"]["apply_contract"] == {
        "patch_file": "patch.json",
        "argv_after_executable": ["apply", str(deck.resolve()), "--patch", "patch.json"],
        "risk_authorization_required": [],
        "argv_after_explicit_risk_authorization": ["apply", str(deck.resolve()), "--patch", "patch.json"],
        "windows_stdin": "do_not_use_powershell_here_string",
    }


def test_focused_inspect_reports_exact_risk_authorization_argv() -> None:
    deck = Path("fixtures/synthetic/synthetic.pptx")
    payload = service.inspect(deck, purpose="content", no_state=True)
    contract = payload["data"]["apply_contract"]

    assert contract["risk_authorization_required"] == ["external_relationship"]
    assert contract["argv_after_executable"][-2:] == ["--patch", "patch.json"]
    assert contract["argv_after_explicit_risk_authorization"][-2:] == ["--allow-risk", "external_relationship"]


def test_inspect_allow_risk_is_compatibility_only_and_security_remains_blocked() -> None:
    deck = Path("fixtures/synthetic/synthetic.pptx")
    payload = service.inspect(
        deck,
        purpose="content",
        no_state=True,
        allow_risk={"external_relationship"},
    )

    assert payload["data"]["ignored_allow_risk"] == ["external_relationship"]
    assert payload["data"]["security"]["safe"] is False
    assert payload["data"]["apply_contract"]["risk_authorization_required"] == ["external_relationship"]


def test_invalid_inspect_fields_return_nearest_candidates() -> None:
    with pytest.raises(PptAgentError) as caught:
        service._project_inspection(
            {"slide_count": 1, "slide_size": [13.33, 7.5], "slides": {}},
            purpose="layout",
            fields=["geomtry"],
        )

    details = caught.value.details
    assert caught.value.code == "INVALID_INSPECT_FIELDS"
    assert details["allowed"] == ["animation", "geometry", "identity", "image", "name", "style", "table", "text", "type"]
    assert details["nearest"] == {"geomtry": ["geometry"]}
    assert details["corrected_example"] == {
        "command": "ppt-agent inspect FILE --for layout --fields geometry,name,text,type",
        "fields": ["geometry", "name", "text", "type"],
    }


def test_accept_requires_matching_review_token(deck: Path) -> None:
    candidate = candidate_path(deck)
    shutil.copy2(deck, candidate)
    current = revision(candidate)
    store = StateStore(document_id(deck))
    store.save_review({
        "candidate_revision": current,
        "review_token_hash": hashlib.sha256(b"correct").hexdigest(),
        "qa_error_count": 0,
        "qa_issues": [],
    })
    with pytest.raises(PptAgentError) as caught:
        service.accept(candidate, current, "wrong")
    assert caught.value.code == "REVIEW_TOKEN_INVALID"
    assert candidate.exists()
    assert not accepted_path(deck).exists()


def test_accept_promotes_candidate_and_discard_removes_it(deck: Path) -> None:
    candidate = candidate_path(deck)
    shutil.copy2(deck, candidate)
    current = revision(candidate)
    store = StateStore(document_id(deck))
    store.save_review({
        "candidate_revision": current,
        "review_token_hash": hashlib.sha256(b"correct").hexdigest(),
        "qa_error_count": 0,
        "qa_issues": [],
    })
    result = service.accept(candidate, current, "correct")
    assert result["ok"] is True
    assert accepted_path(deck).exists()
    assert not candidate.exists()

    shutil.copy2(deck, candidate)
    discarded = service.discard(deck)
    assert discarded["data"]["candidate_removed"] is True
    assert not candidate.exists()


def test_wps_failure_rolls_back_candidate(
    deck: Path, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    original = deck.read_bytes()

    def fake_apply(source: Path, operations: list[dict], output: Path) -> dict:
        shutil.copy2(source, output)
        return {"engine_output": "prepared"}

    def fail_wps(path: Path, operations: list[dict]) -> dict:
        raise PptAgentError("WPS_SAVE_FAILED", "WPS 保存失败", "retry", retryable=True)

    monkeypatch.setattr(service, "apply_operations", fake_apply)
    monkeypatch.setattr(service.wps, "apply_wps_operations", fail_wps)
    patch = _patch(deck, "wps-fail", [
        {"op": "set_text", "object": "s0:s2", "text": "Changed"},
        {"op": "add_animation", "object": "s0:s2", "effect": "fade"},
    ])
    patch_path = tmp_path / "patch.json"
    patch_path.write_text(json.dumps(patch), encoding="utf-8")
    with pytest.raises(PptAgentError) as caught:
        service.apply(deck, patch_path)
    assert caught.value.code == "WPS_SAVE_FAILED"
    assert deck.read_bytes() == original
    assert not candidate_path(deck).exists()
    assert StateStore(document_id(deck)).requests() == {}


def test_animation_schema_accepts_paragraphs_all_and_rejects_legacy_paragraph(deck: Path) -> None:
    from ppt_agent.models import PatchRequest

    valid = PatchRequest.model_validate(_patch(deck, "paragraphs-all", [{
        "op": "add_animation",
        "object": "s0:s2",
        "effect": "appear",
        "trigger": "with_previous",
        "paragraphs": "all",
    }]))
    assert valid.operations[0].paragraphs == "all"
    animation_schema = type(valid.operations[0]).model_json_schema()
    assert "paragraphs" in animation_schema["properties"]
    assert "paragraph" not in animation_schema["properties"]

    with pytest.raises(ValueError):
        PatchRequest.model_validate(_patch(deck, "legacy-paragraph", [{
            "op": "add_animation",
            "object": "s0:s2",
            "paragraph": 1,
        }]))


@pytest.mark.parametrize("animations", [
    [
        {"op": "add_animation", "object": "s0:s2", "paragraphs": "all"},
        {"op": "add_animation", "object": "s0:s2", "effect": "appear", "paragraphs": "all"},
    ],
    [
        {"op": "add_animation", "object": "s0:s2", "paragraph": 1},
        {"op": "add_animation", "object": "s0:s2", "paragraphs": "all"},
    ],
    [
        {"op": "add_animation", "object": "s0:s2"},
        {"op": "add_animation", "object": "s0:s2", "effect": "fade", "trigger": "on_click"},
    ],
    [
        {"op": "add_animation", "object": "s0:s2", "effect": "fade"},
        {"op": "add_animation", "object": "s0:s2", "effect": "appear", "paragraphs": "all"},
    ],
])
def test_duplicate_animation_is_rejected_before_wps_and_without_state_change(
    animations: list[dict], deck: Path, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    original = deck.read_bytes()
    wps_started = False

    def should_not_start(*args, **kwargs):
        nonlocal wps_started
        wps_started = True
        raise AssertionError("WPS must not start for an invalid animation patch")

    monkeypatch.setattr(service.wps, "apply_wps_operations", should_not_start)
    patch_path = tmp_path / "duplicate-animation.json"
    patch_path.write_text(json.dumps(_patch(deck, "duplicate-animation", animations)), encoding="utf-8")

    with pytest.raises(PptAgentError) as caught:
        service.apply(deck, patch_path)

    assert caught.value.code == "DUPLICATE_ANIMATION"
    assert caught.value.next_action == "fix_patch"
    assert caught.value.details == {"object": "s0:s2", "document_unchanged": True}
    assert wps_started is False
    assert deck.read_bytes() == original
    assert not candidate_path(deck).exists()
    assert StateStore(document_id(deck)).requests() == {}


def test_duplicate_animation_cli_error_is_structured_json(
    deck: Path, tmp_path: Path, capsys: pytest.CaptureFixture[str]
) -> None:
    patch_path = tmp_path / "duplicate-animation.json"
    patch_path.write_text(json.dumps(_patch(deck, "duplicate-cli", [
        {"op": "add_animation", "object": "s0:s2"},
        {"op": "add_animation", "object": "s0:s2", "paragraphs": "all"},
    ])), encoding="utf-8")

    with pytest.raises(SystemExit) as caught:
        cli.main(["apply", str(deck), str(patch_path)])

    payload = json.loads(capsys.readouterr().out)
    assert caught.value.code == 1
    assert payload["error_code"] == "DUPLICATE_ANIMATION"
    assert payload["next_action"] == "fix_patch"
    assert payload["document_unchanged"] is True


def test_wps_context_quits_owned_instance_on_failure(monkeypatch: pytest.MonkeyPatch) -> None:
    class App:
        Visible = True
        DisplayAlerts = 1
        quit_called = False

        def Quit(self) -> None:
            self.quit_called = True

    app = App()

    class Client:
        @staticmethod
        def DispatchEx(prog_id: str) -> App:
            assert prog_id == "KWPP.Application"
            return app

    monkeypatch.setattr(wps, "_client", lambda: Client)
    with pytest.raises(PptAgentError) as caught:
        with wps.wps_application():
            raise RuntimeError("operation failed")
    assert caught.value.code == "WPS_COM_FAILED"
    assert app.Visible is False
    assert app.DisplayAlerts == 0
    assert app.quit_called is True


def test_wps_save_failure_closes_presentation_and_quits_app(
    monkeypatch: pytest.MonkeyPatch, tmp_path: Path
) -> None:
    class PresentationHandle:
        closed = False

        def Save(self) -> None:
            raise RuntimeError("save failed")

        def Close(self) -> None:
            self.closed = True

    presentation = PresentationHandle()

    class App:
        Visible = True
        DisplayAlerts = 1
        quit_called = False
        Presentations = type("Presentations", (), {"Open": lambda self, *args: presentation})()

        def Quit(self) -> None:
            self.quit_called = True

    app = App()

    class Client:
        @staticmethod
        def DispatchEx(prog_id: str) -> App:
            return app

    monkeypatch.setattr(wps, "_client", lambda: Client)
    with pytest.raises(PptAgentError) as caught:
        wps.finalize(tmp_path / "deck.pptx")
    assert caught.value.code == "WPS_COM_FAILED"
    assert presentation.closed is True
    assert app.quit_called is True


def test_swap_image_is_translated_and_published(
    deck: Path, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    replacement = tmp_path / "replacement.png"
    replacement.write_bytes(b"local image")
    captured: list[dict] = []

    def fake_apply(source: Path, operations: list[dict], output: Path) -> dict:
        captured.extend(operations)
        shutil.copy2(source, output)
        return {"engine_output": "swapped"}

    monkeypatch.setattr(service, "apply_operations", fake_apply)
    monkeypatch.setattr(service.wps, "finalize", lambda path: {"wps_version": "test"})
    patch = _patch(deck, "swap-image", [{
        "op": "swap_image",
        "object": "s0:s2",
        "image": str(replacement),
    }])
    patch_path = tmp_path / "patch.json"
    patch_path.write_text(json.dumps(patch), encoding="utf-8")
    response = service.apply(deck, patch_path)
    assert response["ok"] is True
    assert captured[0]["op"] == "swap_image"


def test_schema_valid_alternate_forms_apply_first_time_without_internal_error(
    deck: Path, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    # Real engine path (translate_operation included); only WPS is mocked.
    monkeypatch.setattr(service.wps, "finalize", lambda path: {"wps_version": "test"})
    patch = _patch(deck, "alternate-forms", [
        {"op": "resize", "object": "s0:s2", "scale": 0.8},
        {"op": "add_shape", "slide": 0, "kind": "line", "from": [0.0, 0.0], "to": [5.0, 5.0]},
    ])
    patch_path = tmp_path / "alternate-forms.json"
    patch_path.write_text(json.dumps(patch), encoding="utf-8")

    response = service.apply(deck, patch_path)

    assert response["ok"] is True
    assert response["data"]["operations_applied"] == 2
    assert "resize" in response["data"]["engine_output"]
    assert "add-shape" in response["data"]["engine_output"]
