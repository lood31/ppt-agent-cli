from __future__ import annotations

import hashlib
import json
import sys
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def build(output: Path) -> dict[str, object]:
    output.parent.mkdir(parents=True, exist_ok=True)
    image_path = output.parent / "synthetic-image.png"
    image = Image.new("RGB", (900, 500), "#123B4A")
    draw = ImageDraw.Draw(image)
    draw.rectangle((80, 80, 820, 420), outline="#F18F6B", width=12)
    font = ImageFont.truetype("arial.ttf", 42)
    draw.text((160, 215), "SYNTHETIC TEST ASSET", fill="#FFFFFF", font=font)
    image.save(image_path)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(8.5), Inches(0.8))
    p = title.text_frame.paragraphs[0]
    first = p.add_run(); first.text = "WPS "; first.font.bold = True; first.font.size = Pt(32); first.font.color.rgb = RGBColor(20, 56, 70)
    second = p.add_run(); second.text = "保真测试稿"; second.font.size = Pt(32); second.font.color.rgb = RGBColor(225, 91, 61)
    second.hyperlink.address = "https://example.com/ppt-agent"
    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.3), Inches(2.4))
    body.text_frame.clear()
    for index, text in enumerate(["富文本与项目符号", "本地图片与裁剪", "形状、表格与备注"]):
        paragraph = body.text_frame.paragraphs[0] if index == 0 else body.text_frame.add_paragraph()
        paragraph.text = text; paragraph.level = index % 2; paragraph.font.size = Pt(19)
    picture = slide.shapes.add_picture(str(image_path), Inches(6.4), Inches(1.45), width=Inches(5.8), height=Inches(3.2))
    picture.crop_left = 0.05; picture.crop_right = 0.05
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.5), Inches(2.2), Inches(1.1))
    rect.text = "动画目标"; rect.fill.solid(); rect.fill.fore_color.rgb = RGBColor(225, 91, 61)
    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.4), Inches(4.5), Inches(1.1), Inches(1.1))
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(4.5), Inches(5.05), Inches(6.4), Inches(5.05))
    table = slide.shapes.add_table(2, 3, Inches(6.4), Inches(5.0), Inches(5.8), Inches(1.4)).table
    for row, values in enumerate([["指标", "当前", "目标"], ["Token", "100%", "≤30%"]]):
        for col, value in enumerate(values):
            table.cell(row, col).text = value
    slide.notes_slide.notes_text_frame.text = "合成稿，仅用于本机 WPS PoC，不含客户数据。"

    for chart_type, title_text in [
        (XL_CHART_TYPE.COLUMN_CLUSTERED, "柱状图"),
        (XL_CHART_TYPE.LINE, "折线图"),
        (XL_CHART_TYPE.PIE, "饼图"),
    ]:
        chart_slide = prs.slides.add_slide(prs.slide_layouts[6])
        heading = chart_slide.shapes.add_textbox(Inches(0.8), Inches(0.55), Inches(6), Inches(0.7))
        heading.text_frame.text = title_text; heading.text_frame.paragraphs[0].font.size = Pt(30)
        data = ChartData(); data.categories = ["A", "B", "C"]
        data.add_series("演示数据", (3, 5, 4))
        if chart_type == XL_CHART_TYPE.PIE:
            chart = chart_slide.shapes.add_chart(chart_type, Inches(2.7), Inches(1.45), Inches(7.9), Inches(5.35), data).chart
            chart.has_legend = True; chart.legend.position = XL_LEGEND_POSITION.RIGHT
        else:
            chart = chart_slide.shapes.add_chart(chart_type, Inches(1.15), Inches(1.45), Inches(11.0), Inches(5.25), data).chart
            chart.has_legend = chart_type == XL_CHART_TYPE.LINE
            if chart.has_legend:
                chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.has_title = False

    prs.save(output)
    digest = hashlib.sha256(output.read_bytes()).hexdigest()
    return {"file": str(output.resolve()), "sha256": digest, "slides": len(prs.slides), "animation_target_shape_id": rect.shape_id}


if __name__ == "__main__":
    target = Path(sys.argv[1]) if len(sys.argv) > 1 else Path("fixtures/synthetic/synthetic.pptx")
    print(json.dumps(build(target), ensure_ascii=False, indent=2))
