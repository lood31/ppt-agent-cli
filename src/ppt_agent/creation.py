from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt

from .models import CreateSpec, SlideSpec


INK = RGBColor(24, 30, 36)
PAPER = RGBColor(245, 241, 232)
ACCENT = RGBColor(210, 73, 51)
MUTED = RGBColor(102, 111, 119)


def create_presentation(spec: CreateSpec, output: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(spec.width)
    prs.slide_height = Inches(spec.height)
    while prs.slides:
        r_id = prs.slides._sldIdLst[-1].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[-1]
    for index, slide_spec in enumerate(spec.slides):
        _add_slide(prs, slide_spec, index, len(spec.slides))
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


def _add_slide(prs: Presentation, spec: SlideSpec, index: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = INK if spec.layout in {"title", "section", "quote"} else PAPER

    if spec.layout == "blank":
        return
    dark = spec.layout in {"title", "section", "quote"}
    foreground = PAPER if dark else INK
    accent = RGBColor(236, 111, 79) if dark else ACCENT

    if spec.layout in {"title", "section"}:
        _textbox(slide, 0.85, 1.45, 11.5, 2.1, spec.title, 42 if spec.layout == "section" else 48, foreground, True)
        if spec.body:
            _textbox(slide, 0.9, 4.3, 10.4, 1.5, "\n".join(spec.body), 20, foreground)
    elif spec.layout == "quote":
        _textbox(slide, 1.0, 1.2, 0.5, 0.8, "“", 58, accent, True)
        _textbox(slide, 1.45, 1.45, 10.5, 1.0, spec.title, 34, foreground, True)
        if spec.body:
            panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.45), Inches(2.75), Inches(10.6), Inches(3.25))
            panel.fill.solid(); panel.fill.fore_color.rgb = RGBColor(39, 48, 57); panel.line.fill.background()
            _textbox(slide, 1.85, 3.05, 9.8, 2.65, "\n".join(spec.body), 20, foreground)
    elif spec.layout in {"two_column", "comparison"}:
        _title(slide, spec.title, foreground)
        midpoint = max(1, (len(spec.body) + 1) // 2)
        _card(slide, 0.75, 1.75, 5.8, 4.75, spec.body[:midpoint], foreground, accent)
        _card(slide, 6.8, 1.75, 5.8, 4.75, spec.body[midpoint:], foreground, accent)
    elif spec.layout == "stat":
        _title(slide, spec.title, foreground)
        for item_index, text in enumerate(spec.body[:3]):
            x = 0.8 + item_index * 4.15
            card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(1.9), Inches(3.65), Inches(3.7))
            card.fill.solid(); card.fill.fore_color.rgb = RGBColor(255, 253, 248); card.line.color.rgb = RGBColor(220, 214, 202)
            _textbox(slide, x + 0.3, 2.35, 3.05, 2.8, text, 27, accent, True, align=PP_ALIGN.CENTER, vertical=MSO_VERTICAL_ANCHOR.MIDDLE)
    elif spec.layout == "timeline":
        _title(slide, spec.title, foreground)
        for item_index, text in enumerate(spec.body[:4]):
            x = 0.8 + item_index * 3.1
            circle = slide.shapes.add_shape(9, Inches(x), Inches(2.35), Inches(0.55), Inches(0.55))
            circle.fill.solid(); circle.fill.fore_color.rgb = accent; circle.line.fill.background()
            _textbox(slide, x, 2.46, 0.55, 0.3, str(item_index + 1), 13, PAPER, True, align=PP_ALIGN.CENTER)
            if item_index < min(3, len(spec.body) - 1):
                connector = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(x + 0.85), Inches(2.43), Inches(1.65), Inches(0.36))
                connector.fill.solid(); connector.fill.fore_color.rgb = RGBColor(230, 196, 106); connector.line.fill.background()
            _textbox(slide, x - 0.15, 3.15, 2.75, 2.0, text, 17, foreground, item_index == 0)
    elif spec.layout == "image_text":
        _title(slide, spec.title, foreground)
        block = slide.shapes.add_shape(1, Inches(0.75), Inches(1.7), Inches(5.1), Inches(4.95))
        block.fill.solid(); block.fill.fore_color.rgb = accent; block.line.fill.background()
        _textbox(slide, 6.35, 1.85, 5.8, 4.5, "\n".join(spec.body), 21, foreground)
    else:
        _title(slide, spec.title, foreground)
        _textbox(slide, 0.9, 1.75, 11.5, 4.7, "\n".join(f"— {item}" for item in spec.body), 22, foreground)

    _textbox(slide, 0.85, 6.78, 4.0, 0.35, f"PPT AGENT  /  {index + 1:02d}", 11, MUTED)
    if spec.notes:
        notes_frame = slide.notes_slide.notes_text_frame
        notes_frame.text = spec.notes


def _title(slide, text: str, color: RGBColor) -> None:
    _textbox(slide, 0.75, 0.55, 11.8, 0.75, text, 34, color, True)


def _textbox(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    size: float,
    color: RGBColor,
    bold: bool = False,
    *,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    vertical: MSO_VERTICAL_ANCHOR = MSO_VERTICAL_ANCHOR.TOP,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = vertical
    lines = text.splitlines() or [""]
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color
        paragraph.alignment = align
        paragraph.space_after = Pt(8)
    return shape


def _card(slide, x: float, y: float, w: float, h: float, lines: list[str], color: RGBColor, accent: RGBColor) -> None:
    card = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid(); card.fill.fore_color.rgb = RGBColor(255, 253, 248); card.line.color.rgb = RGBColor(220, 214, 202)
    if lines:
        _textbox(slide, x + 0.35, y + 0.35, w - 0.7, h - 0.7, "\n".join(f"— {item}" for item in lines), 18, color)
