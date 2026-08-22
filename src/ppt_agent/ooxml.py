from __future__ import annotations

import os
import re
import time
import zipfile
from collections import Counter
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn

from .errors import PptAgentError


RISK_PARTS = {
    "macro": ("vbaproject.bin",),
    "activex": ("activex/",),
    # Native chart workbooks also live under ppt/embeddings and are expected.
    # Only explicit OLE object parts/relationships are high risk.
    "ole": ("oleobject",),
}

MC_NS = "http://schemas.openxmlformats.org/markup-compatibility/2006"
_TRANSITION_EL_TO_TYPE = {
    "fade": "fade",
    "cut": "cut",
    "dissolve": "dissolve",
    "push": "push",
    "wipe": "wipe",
    "split": "split",
    "cover": "cover",
    "pull": "uncover",
    "zoom": "zoom",
}


def validate_pptx(path: Path) -> None:
    if not path.exists() or not path.is_file():
        raise PptAgentError("FILE_NOT_FOUND", "PPTX 文件不存在", "check_path")
    if path.suffix.lower() != ".pptx":
        raise PptAgentError("INVALID_FILE_TYPE", "仅支持 .pptx 文件", "use_pptx")
    if not zipfile.is_zipfile(path):
        raise PptAgentError("INVALID_PPTX", "文件不是有效的 PPTX ZIP 容器", "repair_or_replace")
    try:
        with zipfile.ZipFile(path) as archive:
            names = set(archive.namelist())
            if "[Content_Types].xml" not in names or "ppt/presentation.xml" not in names:
                raise PptAgentError("INVALID_PPTX", "PPTX 缺少必要的 OOXML 部件", "repair_or_replace")
            bad = archive.testzip()
            if bad:
                raise PptAgentError("CORRUPT_PPTX", f"PPTX 部件损坏：{bad}", "repair_or_replace")
    except zipfile.BadZipFile as exc:
        raise PptAgentError("INVALID_PPTX", "PPTX ZIP 结构无法读取", "repair_or_replace") from exc


def security_scan(path: Path) -> dict[str, Any]:
    validate_pptx(path)
    risks: list[dict[str, Any]] = []
    with zipfile.ZipFile(path) as archive:
        names = archive.namelist()
        lowered = [name.lower() for name in names]
        for kind, markers in RISK_PARTS.items():
            matches = [name for name, low in zip(names, lowered) if any(marker in low for marker in markers)]
            if matches:
                risks.append({"type": kind, "severity": "error", "count": len(matches)})

        external_count = 0
        remote_media_count = 0
        for name in names:
            if not name.endswith(".rels"):
                continue
            try:
                root = ET.fromstring(archive.read(name))
            except ET.ParseError:
                continue
            for rel in root:
                if rel.attrib.get("TargetMode") != "External":
                    continue
                external_count += 1
                target = rel.attrib.get("Target", "")
                rel_type = rel.attrib.get("Type", "")
                if re.match(r"https?://", target, re.I) and rel_type.endswith(("/image", "/video", "/audio", "/media")):
                    remote_media_count += 1
        if external_count:
            risks.append({"type": "external_relationship", "severity": "error", "count": external_count})
        if remote_media_count:
            risks.append({"type": "remote_media", "severity": "error", "count": remote_media_count})

        counts = Counter()
        for name in lowered:
            if name.startswith("ppt/slides/slide") and name.endswith(".xml"):
                counts["slides"] += 1
            elif name.startswith("ppt/media/"):
                counts["media"] += 1
            elif name.startswith("ppt/charts/chart") and name.endswith(".xml"):
                counts["charts"] += 1

    return {
        "safe": not any(item["severity"] == "error" for item in risks),
        "risks": risks,
        "counts": dict(counts),
    }


def enforce_risk_policy(path: Path, allow_risk: set[str]) -> dict[str, Any]:
    report = security_scan(path)
    blocked = [risk for risk in report["risks"] if risk["type"] not in allow_risk]
    if blocked:
        raise PptAgentError(
            "SECURITY_RISK_BLOCKED",
            "文档包含默认阻断的高风险特性",
            "inspect_risks_or_allow_explicitly",
            details={"blocked_risks": blocked, "document_unchanged": True},
        )
    return report


# ---------------------------------------------------------------------------
# Slide transitions.
#
# WPS wraps transitions it writes in mc:AlternateContent (Choice for the p14
# duration, Fallback for the plain element). The pinned engine only looks for a
# direct p:transition child, so these helpers provide one authoritative parse
# for inspect, diff, and the set_slide pre-clean step.
# ---------------------------------------------------------------------------

def _transition_element(sld: Any) -> Any:
    direct = sld.find(qn("p:transition"))
    if direct is not None:
        return direct
    for alternate in sld.findall("{%s}AlternateContent" % MC_NS):
        fallback = alternate.find("{%s}Fallback" % MC_NS)
        if fallback is not None:
            element = fallback.find(qn("p:transition"))
            if element is not None:
                return element
    return None


def _parse_transition_element(element: Any) -> dict[str, Any]:
    out: dict[str, Any] = {}
    for child in element:
        local = child.tag.rsplit("}", 1)[-1]
        out["type"] = _TRANSITION_EL_TO_TYPE.get(local, local)
        for key in ("dir", "orient"):
            if child.get(key):
                out[key] = child.get(key)
        break
    if element.get("spd"):
        out["speed"] = element.get("spd")
    if element.get("advTm") is not None:
        out["advance_after"] = int(element.get("advTm")) / 1000.0
    if element.get("advClick") == "0":
        out["advance_on_click"] = False
    return out


def transitions_by_slide(path: Path) -> dict[int, dict[str, Any] | None]:
    """Authoritative per-slide transition facts, direct or mc-wrapped."""
    presentation = Presentation(path)
    result: dict[int, dict[str, Any] | None] = {}
    for index, slide in enumerate(presentation.slides):
        element = _transition_element(slide._element)
        result[index] = _parse_transition_element(element) if element is not None else None
    return result


def strip_transition_elements(path: Path, slides: set[int] | None = None) -> None:
    """Remove direct and mc-wrapped p:transition from selected slides in place."""
    presentation = Presentation(path)
    for index, slide in enumerate(presentation.slides):
        if slides is not None and index not in slides:
            continue
        sld = slide._element
        direct = sld.find(qn("p:transition"))
        if direct is not None:
            sld.remove(direct)
        for alternate in sld.findall("{%s}AlternateContent" % MC_NS):
            if alternate.find(".//" + qn("p:transition")) is not None:
                sld.remove(alternate)
    presentation.save(path)


_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _slide_part_paths(path: Path) -> list[str]:
    """Slide index (document order) -> zip part path."""
    with zipfile.ZipFile(path) as archive:
        rels_root = ET.fromstring(archive.read("ppt/_rels/presentation.xml.rels"))
        rid_to_target: dict[str, str] = {}
        for rel in rels_root:
            rid = rel.attrib.get("Id")
            target = rel.attrib.get("Target", "")
            if rid and target:
                rid_to_target[rid] = target
        pres_root = ET.fromstring(archive.read("ppt/presentation.xml"))
        parts: list[str] = []
        for sld_id in pres_root.findall(".//{%s}sldId" % _P_NS):
            rid = sld_id.attrib.get("{%s}id" % _R_NS, "")
            target = rid_to_target.get(rid, "")
            if not target:
                raise PptAgentError("TRANSITION_PARSE_FAILED", "无法解析幻灯片部件关系", "report_bug")
            parts.append(target.lstrip("/") if target.startswith("/") else "ppt/" + target)
        if not parts:
            raise PptAgentError("TRANSITION_PARSE_FAILED", "演示文稿没有幻灯片部件", "report_bug")
        return parts


def _restore_slide_transitions(xml_bytes: bytes, spec: dict[str, Any]) -> bytes:
    """Rewrite one slide XML: re-apply spd and effect dir/orient, verify type."""
    root = etree.fromstring(xml_bytes)
    transition_elements: list[Any] = []
    for child in root:
        local = child.tag.rsplit("}", 1)[-1]
        if local == "transition" and child.tag == "{%s}transition" % _P_NS:
            transition_elements.append(child)
        elif local == "AlternateContent":
            transition_elements.extend(child.iter("{%s}transition" % _P_NS))
    if not transition_elements:
        raise PptAgentError(
            "TRANSITION_FIDELITY_FAILED",
            "目标页未找到切换效果，无法恢复属性",
            "reinspect",
            details={"document_unchanged": True},
        )
    expected_type = spec.get("type")
    restored = False
    for element in transition_elements:
        effect = None
        actual_type: str | None = None
        for child in element:
            local = child.tag.rsplit("}", 1)[-1]
            canonical = _TRANSITION_EL_TO_TYPE.get(local)
            if canonical is None:
                continue
            effect = child
            actual_type = canonical
            break
        if effect is None or actual_type is None:
            raise PptAgentError(
                "TRANSITION_FIDELITY_FAILED",
                "切换元素缺少效果子节点",
                "reinspect",
                details={"document_unchanged": True},
            )
        if expected_type is not None and actual_type != expected_type:
            raise PptAgentError(
                "TRANSITION_FIDELITY_FAILED",
                "WPS 保存后的切换类型与预期不一致",
                "reinspect",
                details={"expected": expected_type, "actual": actual_type, "document_unchanged": True},
            )
        if "speed" in spec:
            element.set("spd", spec["speed"])
        for key in ("dir", "orient"):
            if spec.get(key) is not None:
                effect.set(key, spec[key])
        restored = True
    if not restored:
        raise PptAgentError("TRANSITION_FIDELITY_FAILED", "切换属性恢复失败", "report_bug")
    return etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)


def _replace_with_retry(source: Path, target: Path, attempts: int = 5) -> None:
    """os.replace 的短有界重试：Windows 下瞬时占用（WinError 5/32）常见于
    反病毒/索引器短暂打开刚写出的文件，退避重试即可，不改变失败语义。"""
    delay = 0.02
    for attempt in range(attempts):
        try:
            os.replace(source, target)
            return
        except OSError as exc:
            if getattr(exc, "winerror", None) in (5, 32) and attempt < attempts - 1:
                time.sleep(delay)
                delay *= 2
                continue
            raise


def restore_transition_options(path: Path, expected: dict[int, dict[str, Any]]) -> None:
    """Surgically re-apply dir/orient/speed that WPS normalizes away on save.

    Only the affected slide XML parts are rewritten; every other ZIP entry
    keeps its original bytes. Fails with TRANSITION_FIDELITY_FAILED when the
    transition WPS actually saved does not match the requested type.
    """
    if not expected:
        return
    try:
        part_paths = _slide_part_paths(path)
        with zipfile.ZipFile(path) as archive:
            infos = archive.infolist()
            entries = {info.filename: archive.read(info.filename) for info in infos}
        for index, spec in expected.items():
            if index >= len(part_paths):
                raise PptAgentError(
                    "TRANSITION_FIDELITY_FAILED",
                    "切换目标页超出范围",
                    "reinspect",
                    details={"slide": index, "document_unchanged": True},
                )
            part = part_paths[index]
            if part not in entries:
                raise PptAgentError("TRANSITION_FIDELITY_FAILED", "未找到目标页部件", "report_bug")
            entries[part] = _restore_slide_transitions(entries[part], spec)
        temp = path.with_name(f".{path.name}.transition.tmp")
        try:
            with zipfile.ZipFile(temp, "w", zipfile.ZIP_DEFLATED) as out:
                for info in infos:
                    out.writestr(info, entries[info.filename])
            _replace_with_retry(temp, path)
        finally:
            try:
                temp.unlink()
            except FileNotFoundError:
                pass
    except PptAgentError:
        raise
    except Exception as exc:
        raise PptAgentError(
            "TRANSITION_PARSE_FAILED",
            "切换属性恢复失败",
            "report_bug",
            details={"error": str(exc)[:500]},
        ) from exc
