from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation

from .engine import inspect_json


def run_qa(path: Path, profile: str = "basic", *, suggest_fixes: bool = False) -> dict[str, Any]:
    presentation = Presentation(path)
    width = presentation.slide_width
    height = presentation.slide_height
    issues: list[dict[str, Any]] = []
    suggested_operations: list[dict[str, Any]] = []

    for slide_index, slide in enumerate(presentation.slides):
        for shape in slide.shapes:
            if shape.left < 0 or shape.top < 0 or shape.left + shape.width > width or shape.top + shape.height > height:
                issues.append({
                    "severity": "error",
                    "code": "OUT_OF_BOUNDS",
                    "slide": slide_index,
                    "object": f"s{slide_index}:s{shape.shape_id}",
                })
                if suggest_fixes and shape.width <= width and shape.height <= height:
                    suggested_operations.append({
                        "op": "move",
                        "object": f"s{slide_index}:s{shape.shape_id}",
                        "x": round(max(0, min(shape.left, width - shape.width)) / 914400, 3),
                        "y": round(max(0, min(shape.top, height - shape.height)) / 914400, 3),
                    })
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text.strip()
            if not text:
                continue
            if profile in {"presentation", "assignment"}:
                sizes = [run.font.size.pt for paragraph in shape.text_frame.paragraphs for run in paragraph.runs if run.font.size]
                minimum = min(sizes) if sizes else None
                floor = 12 if profile == "assignment" else 14
                if minimum is not None and minimum < floor:
                    issues.append({
                        "severity": "warning",
                        "code": "FONT_TOO_SMALL",
                        "slide": slide_index,
                        "object": f"s{slide_index}:s{shape.shape_id}",
                        "value": minimum,
                        "minimum": floor,
                    })
                    if suggest_fixes:
                        suggested_operations.append({
                            "op": "set_style",
                            "object": f"s{slide_index}:s{shape.shape_id}",
                            "font_size": floor,
                        })
            if len(text) > 800:
                issues.append({
                    "severity": "warning",
                    "code": "TEXT_TOO_DENSE",
                    "slide": slide_index,
                    "object": f"s{slide_index}:s{shape.shape_id}",
                })

    engine_data = inspect_json(path)
    for slide_index, shapes in engine_data.get("slides", {}).items():
        for shape_id, shape in shapes.items():
            if shape_id.startswith("_") or not isinstance(shape, dict):
                continue
            raw_issues = shape.get("issues", [])
            normalized: list[str] = []
            if isinstance(raw_issues, dict):
                for key, value in raw_issues.items():
                    if key == "warnings" and isinstance(value, list):
                        normalized.extend(str(item) for item in value)
                    else:
                        normalized.append(f"{key}:{value}")
            elif isinstance(raw_issues, list):
                normalized.extend(str(item) for item in raw_issues)
            for issue in normalized:
                issues.append({
                    "severity": "warning",
                    "code": "ENGINE_LAYOUT_WARNING",
                    "slide": int(slide_index),
                    "object": f"s{slide_index}:{shape_id}",
                    "detail": str(issue)[:300],
                })

    report = {
        "profile": profile,
        "issue_count": len(issues),
        "error_count": sum(item["severity"] == "error" for item in issues),
        "warning_count": sum(item["severity"] == "warning" for item in issues),
        "issues": issues,
    }
    if suggest_fixes:
        report["suggested_patch"] = {"operations": suggested_operations} if suggested_operations else None
        report["suggested_fix_count"] = len(suggested_operations)
    return report
